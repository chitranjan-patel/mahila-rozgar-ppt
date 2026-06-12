import os
from pptx.util import Inches

ppt_script = r'C:\Users\Chitranjan Kumar\Desktop\ppt\create_ppt.py'

with open(ppt_script, 'r', encoding='utf-8') as f:
    content = f.read()

# I will simply replace `add_slide(prs, title_text, content_items, image_desc, subtitle=None, is_title_slide=False):`
# with a version that handles `image_path` as well.
# Actually, I can just replace `image_desc` with `image_path` everywhere in main.

new_func = """def add_slide(prs, title_text, content_items, image_path, subtitle=None, is_title_slide=False):
    layout = prs.slide_layouts[0] if is_title_slide else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(0, 51, 153)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER if is_title_slide else PP_ALIGN.LEFT

    if is_title_slide:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"{subtitle}"
        if image_path and os.path.exists(image_path):
            # Center the image on title slide
            from pptx.util import Inches
            slide.shapes.add_picture(image_path, Inches(3), Inches(4.5), width=Inches(4), height=Inches(2.5))
        return slide
    
    # Body
    body_shape = slide.placeholders[1]
    # Reduce body shape width to make room for image
    body_shape.width = Pt(400)
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.size = Pt(20)

    if image_path and os.path.exists(image_path):
        from pptx.util import Inches
        # Add image on the right side
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.5), width=Inches(4.5), height=Inches(5.5))

    return slide
"""

# Replace the function definition
import re
content = re.sub(r'def add_slide\(.*?:.*?(?=def main\(\):)', new_func + '\n', content, flags=re.DOTALL)

# Now in main(), replace `add_slide(prs, "...", [...], "image desc", ...)`
# We need to replace the 4th argument ("image desc") with `f"web/images/slide{slide_number}.png"`

import ast
# To do this safely via regex without ast parsing is tricky because of brackets, but let's just count slides 1 to 15.
# Every add_slide call ends with `], "some string" ... )`
# Let's write a targeted replace for the image description strings.

count = 1
def replace_img_desc(m):
    global count
    # m.group(1) is `], `
    # m.group(2) is the quote
    res = m.group(1) + f"r'C:\\Users\\Chitranjan Kumar\\Desktop\\ppt\\web\\images\\slide{count}.png'"
    count += 1
    return res

content = re.sub(r'(\],\s*)(["\'][^"\']+["\'])', replace_img_desc, content)

with open(ppt_script, 'w', encoding='utf-8') as f:
    f.write("import os\n" + content)

print("Updated create_ppt.py")
