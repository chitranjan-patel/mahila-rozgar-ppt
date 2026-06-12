import os
from pptx import Presentation
from pptx.util import Inches

def main():
    prs = Presentation()
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # typically 6 is blank

    for i in range(1, 16):
        img_path = f'web/images/slide{i}.png'
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_layout)
            # Add image to cover full slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    output_path = "AI_Digital_Marketing_Mahila_Rozgar.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully at: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    main()
